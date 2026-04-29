"""
Generate a PowerPoint slide with the compiled results table and a chart.
Creates: outputs/compiled_results_table.pptx and outputs/compiled_results_chart.png

Usage:
  conda activate mlma
  python scripts/generate_results_ppt.py

Requires: pandas, matplotlib, python-pptx
Install if needed: pip install pandas matplotlib python-pptx
"""

from pathlib import Path
import sys

try:
    import pandas as pd
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    print("Missing dependency:", e)
    print("Install required packages: pip install pandas matplotlib python-pptx")
    sys.exit(1)

BASE = Path(__file__).resolve().parents[1]
CSV_PATH = BASE / "outputs" / "compiled_results_table.csv"
OUT_PPT = BASE / "outputs" / "compiled_results_table.pptx"
OUT_IMG = BASE / "outputs" / "compiled_results_chart.png"

if not CSV_PATH.exists():
    raise FileNotFoundError(f"CSV not found: {CSV_PATH}")

# Read table and keep first five columns (as requested earlier)
df = pd.read_csv(CSV_PATH).iloc[:, :5]
# Replace NaN with empty string for table display
df_display = df.fillna("")

# Create a simple bar chart of accuracy and f1_stress where available
plt.rcParams.update({"figure.autolayout": True})
fig, ax = plt.subplots(figsize=(8, 4))
models = df.index.tolist()
if df.index.dtype == "int64" or df.index.dtype == "float64":
    # reset index names if numeric
    models = df.index.astype(str).tolist()
else:
    models = df["label"].astype(str).tolist() if "label" in df.columns else df.index.astype(str).tolist()

# try to parse numeric columns
acc = pd.to_numeric(df["accuracy"], errors="coerce") if "accuracy" in df.columns else None
f1 = pd.to_numeric(df["f1_stress"], errors="coerce") if "f1_stress" in df.columns else None

x = range(len(models))
if acc is not None:
    ax.bar([i - 0.15 for i in x], acc.fillna(0), width=0.3, label="accuracy")
if f1 is not None:
    ax.bar([i + 0.15 for i in x], f1.fillna(0), width=0.3, label="f1_stress")

ax.set_xticks(x)
ax.set_xticklabels(models, rotation=45, ha="right")
ax.set_ylabel("Score")
ax.set_ylim(0, 1)
ax.legend()
plt.title("Compiled Results: accuracy & f1_stress")

# Save chart image
OUT_IMG.parent.mkdir(parents=True, exist_ok=True)
fig.savefig(OUT_IMG, dpi=150)
plt.close(fig)

# Build PPT
prs = Presentation()
blank_layout = prs.slide_layouts[6]  # blank layout
slide = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(0.2)
width = Inches(9)
height = Inches(0.6)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Compiled Results Summary"
run.font.size = Pt(24)

# Add chart image
img_left = Inches(0.5)
img_top = Inches(1.0)
img_width = Inches(9)
slide.shapes.add_picture(str(OUT_IMG), img_left, img_top, width=img_width)

# Add table below image
rows = len(df_display) + 1
cols = min(5, len(df_display.columns))
if rows > 0:
    table_left = Inches(0.5)
    table_top = Inches(4.6)
    table_width = Inches(9)
    table_height = Inches(1.2 + 0.2 * len(df_display))
    table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

    # set column headers
    for c_idx, col in enumerate(df_display.columns[:cols]):
        table.cell(0, c_idx).text = str(col)
        table.cell(0, c_idx).text_frame.paragraphs[0].runs[0].font.bold = True

    # populate rows
    for r_idx, (_, row) in enumerate(df_display.iterrows(), start=1):
        for c_idx, col in enumerate(df_display.columns[:cols]):
            val = row[col]
            table.cell(r_idx, c_idx).text = f"{val}" if (val is not None and val != "") else ""

# Save PPT
OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPT)
print(f"Wrote {OUT_PPT}")
print(f"Chart image: {OUT_IMG}")
